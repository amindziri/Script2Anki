import os
from pptx import Presentation
from docx import Document
import fitz 


def extract_text_from_pdf(pdf_path):
    text = ""
    doc = fitz.open(pdf_path)
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_pptx(pptx_path):
    text = ""
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_docx(docx_path):
    text = ""
    doc = Document(docx_path)
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_file(file_path):
    # Bestimme den Dateityp basierend auf der Dateierweiterung
    _, file_extension = os.path.splitext(file_path)

    if file_extension.lower() == '.pdf':
        return extract_text_from_pdf(file_path)
    elif file_extension.lower() == '.pptx':
        return extract_text_from_pptx(file_path)
    elif file_extension.lower() == '.docx':
        return extract_text_from_docx(file_path)
    else:
        raise ValueError("Unsupported file format")

# Beispielverwendung:
file_path = "C:/Users/maher/OneDrive/Desktop/test.txt"  # Pfad zur Datei hier einfügen
try:
    extracted_text = extract_text_from_file(file_path)
    print(extracted_text)
except ValueError as e:
    print(e)